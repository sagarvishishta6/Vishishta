from pptx import Presentation

# Load the presentation
prs = Presentation(r"c:\Users\sagar\Desktop\data_YourNextPurchase\data\Hackathon Centrale-ESSEC - Opening presentation.pptx")

# Save to text file to avoid encoding issues
with open('ppt_content.txt', 'w', encoding='utf-8') as f:
    f.write(f"Total slides: {len(prs.slides)}\n\n")
    
    # Extract text from all slides
    for slide_num, slide in enumerate(prs.slides, 1):
        f.write(f"\n{'='*80}\n")
        f.write(f"SLIDE {slide_num}\n")
        f.write(f"{'='*80}\n\n")
        
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                if shape.text.strip():
                    f.write(shape.text + "\n\n")

print("Content extracted to ppt_content.txt")
